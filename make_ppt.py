from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BG_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD  = RGBColor(0x16, 0x21, 0x3E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x99, 0x99, 0x99)

P_INIT   = RGBColor(0x00, 0x7B, 0xFF)
P_WSS    = RGBColor(0x00, 0xC9, 0x7B)
P_SPEED  = RGBColor(0xFF, 0x8C, 0x00)
P_BRAKE  = RGBColor(0xFF, 0x35, 0x5E)
P_PARK   = RGBColor(0xA8, 0x5C, 0xFF)
P_FIN    = RGBColor(0x77, 0x88, 0x99)

# 시스템 색상
S_ECU    = RGBColor(0x00, 0xD2, 0xFF)  # 시안
S_PLC    = RGBColor(0xFF, 0xD6, 0x00)  # 골드
S_NI     = RGBColor(0xFF, 0x69, 0xB4)  # 핑크

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_box(slide, left, top, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=10, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def box_text(shape, text, size=10, color=C_WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tf

def add_sys_badges(slide, x, y, systems):
    """ECU/PLC/NI 배지 그리기"""
    badge_w = Inches(0.32)
    badge_h = Pt(12)
    gap = Inches(0.02)
    for i, (label, clr) in enumerate(systems):
        bx = x + i * (badge_w + gap)
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, y, badge_w, badge_h)
        b.fill.solid()
        b.fill.fore_color.rgb = clr
        b.line.fill.background()
        tf = b.text_frame
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(6)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

def get_sys_badges(sys_str):
    """'EPN' → [('ECU', S_ECU), ('PLC', S_PLC), ('NI', S_NI)]"""
    result = []
    if 'E' in sys_str: result.append(('ECU', S_ECU))
    if 'P' in sys_str: result.append(('PLC', S_PLC))
    if 'N' in sys_str: result.append(('NI', S_NI))
    return result

# ================================================================
# Slide 1: Title
# ================================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s1, BG_DARK)

add_rect(s1, Inches(3), Inches(2.8), Inches(10), Pt(3), P_INIT)
add_text(s1, Inches(3), Inches(3.0), Inches(10), Inches(1),
         "KI-RnB ABS/ESC", 44, C_WHITE, True)
add_text(s1, Inches(3), Inches(3.7), Inches(10), Inches(0.8),
         "Roller Bench Test Sequence", 32, RGBColor(0x88, 0xBB, 0xEE))
add_text(s1, Inches(3), Inches(4.5), Inches(10), Inches(0.5),
         "LX3 Driving Curve  |  33 Steps  |  167 sec  |  Max 100 km/h", 16, C_GRAY)

# Phase preview
phases_preview = [
    ("ECU INIT", P_INIT), ("DTC", RGBColor(0x00,0xB4,0xD8)), ("WSS", P_WSS),
    ("SPEED", P_SPEED), ("BRAKE/ABS", P_BRAKE), ("PARKING", P_PARK), ("FINISH", P_FIN)
]
for i, (name, clr) in enumerate(phases_preview):
    box = add_box(s1, Inches(3) + i * Inches(1.42), Inches(5.5), Inches(1.3), Inches(0.4), clr)
    box_text(box, name, 9, C_WHITE, True)

# System legend
add_text(s1, Inches(3), Inches(6.3), Inches(3), Inches(0.3), "System Legend:", 11, C_GRAY, True)
for i, (label, clr) in enumerate([("ECU  Diagnostic Communication", S_ECU),
                                    ("PLC  Machine & Roller Control", S_PLC),
                                    ("NI   Sensor Data Acquisition", S_NI)]):
    bx = Inches(3)
    by = Inches(6.65) + i * Inches(0.3)
    b = add_box(s1, bx, by, Inches(0.45), Inches(0.22), clr)
    box_text(b, label.split()[0], 8, RGBColor(0,0,0), True)
    add_text(s1, bx + Inches(0.55), by, Inches(3), Inches(0.22),
             " ".join(label.split()[1:]), 10, RGBColor(0xBB, 0xBB, 0xBB))

# ================================================================
# Slide 2: Process Flow (with system badges)
# ================================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s2, BG_DARK)

add_text(s2, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), "PROCESS FLOW", 24, C_WHITE, True)
add_text(s2, Inches(0.5), Inches(0.7), Inches(8), Inches(0.3),
         "Full sequence: 33 steps in 6 phases", 11, C_GRAY)

# System legend (top right)
for i, (label, clr) in enumerate([("ECU", S_ECU), ("PLC", S_PLC), ("NI", S_NI)]):
    lx = Inches(13.5) + i * Inches(0.7)
    lb = add_box(s2, lx, Inches(0.4), Inches(0.55), Inches(0.25), clr)
    box_text(lb, label, 8, RGBColor(0,0,0), True)

# num, name, dur, spd, systems (E=ECU, P=PLC, N=NI)
phases = [
    ("PHASE 1\nECU INITIALIZE", "0~26s", P_INIT, [
        ("01","INITIALIZATION","3s","0","P"),
        ("02","ECU CONNECT","3s","0","E"),
        ("03","START COMM","5s","0","E"),
        ("04","ECU IDENT","3s","0","E"),
        ("05","CHECK SIGNAL","3s","0","E"),
        ("06","DTC READ","3s","0","E"),
        ("07","DTC CLEAR","3s","0","E"),
        ("08","[N] STOP","3s","0","P"),
    ]),
    ("PHASE 2\nWSS TEST", "26~46s", P_WSS, [
        ("09","WSS READY","4s","5","PN"),
        ("10","WSS TEST","10s","5","EPN"),
        ("11","[N] STOP","3s","0","P"),
        ("12","MODE CHECK","3s","0","E"),
    ]),
    ("PHASE 3\nSPEED & CRUISE", "46~101s", P_SPEED, [
        ("13","SPEED UP","10s","40","PN"),
        ("14","SMT TEST","5s","40","EPN"),
        ("15","CRUISE UP","10s","80","PN"),
        ("16","CRUISE CHECK","20s","80","EPN"),
        ("17","MAX SPEED","10s","100","PN"),
    ]),
    ("PHASE 4\nBRAKE & ABS", "101~126s", P_BRAKE, [
        ("18","ENGINE BRAKE","10s","90","PN"),
        ("19","DRAG TEST","3s","80","PN"),
        ("20","BRAKE TEST","3s","60","PN"),
        ("21","DYNAMIC READY","2s","55","PN"),
        ("22","DYNAMIC TEST","4s","10","EPN"),
        ("23","[N] STOP","3s","0","P"),
    ]),
    ("PHASE 5\nREVERSE & PARKING", "126~149s", P_PARK, [
        ("24","REVERSE [R]","5s","0","P"),
        ("25","[R] 5KPH","5s","5","PN"),
        ("26","[N] STOP","3s","0","P"),
        ("27","PARKING READY","5s","0","P"),
        ("28","PARKING TEST","5s","0","PN"),
    ]),
    ("PHASE 6\nFINISH", "149~167s", P_FIN, [
        ("29","IDLE","3s","0","E"),
        ("30","DTC READ","4s","0","E"),
        ("31","DTC CLEAR","3s","0","E"),
        ("32","ECU DISCONNECT","3s","0","E"),
        ("33","FINISH","5s","0","P"),
    ]),
]

col_w = Inches(2.38)
col_gap = Inches(0.12)
start_x = Inches(0.4)
header_y = Inches(1.2)

for pi, (title, trange, color, steps) in enumerate(phases):
    x = start_x + pi * (col_w + col_gap)

    hdr = add_box(s2, x, header_y, col_w, Inches(0.7), color)
    tf = hdr.text_frame
    tf.margin_left = Pt(8)
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = trange
    r2.font.size = Pt(8)
    r2.font.color.rgb = C_WHITE

    add_rect(s2, x + col_w // 2 - Pt(1), header_y + Inches(0.7), Pt(2), Inches(0.12), color)

    step_y = header_y + Inches(0.85)
    step_h = Inches(0.52)
    step_gap = Inches(0.08)

    for si, (num, name, dur, spd, sys) in enumerate(steps):
        sy = step_y + si * (step_h + step_gap)
        card = add_box(s2, x, sy, col_w, step_h, BG_CARD, color)

        csz = Inches(0.3)
        circ = s2.shapes.add_shape(MSO_SHAPE.OVAL,
            x + Inches(0.08), sy + (step_h - csz) // 2, csz, csz)
        circ.fill.solid()
        circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        box_text(circ, num, 8, C_WHITE, True)

        add_text(s2, x + Inches(0.45), sy + Inches(0.02),
                 Inches(1.3), Inches(0.22), name, 9, C_WHITE, True)

        info = f"{dur}  |  {spd} km/h" if spd != "0" else dur
        add_text(s2, x + Inches(0.45), sy + Inches(0.24),
                 Inches(1.0), Inches(0.18), info, 7, C_GRAY)

        # System badges (right side of card)
        badges = get_sys_badges(sys)
        badge_x = x + col_w - Inches(0.12) - len(badges) * Inches(0.34)
        add_sys_badges(s2, badge_x, sy + Inches(0.3), badges)

    if pi < len(phases) - 1:
        arr = s2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            x + col_w + Pt(1), header_y + Inches(0.3), Inches(0.1), Inches(0.15))
        arr.fill.solid()
        arr.fill.fore_color.rgb = C_GRAY
        arr.line.fill.background()

# ================================================================
# Slide 3: Speed Profile (unchanged)
# ================================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s3, BG_DARK)

add_text(s3, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), "SPEED PROFILE", 24, C_WHITE, True)
add_text(s3, Inches(0.5), Inches(0.7), Inches(8), Inches(0.3),
         "Vehicle speed over time (167 seconds)", 11, C_GRAY)

ch_l = Inches(1.3)
ch_t = Inches(1.5)
ch_w = Inches(14.0)
ch_h = Inches(5.8)
max_t = 167
max_s = 110

add_box(s3, ch_l - Inches(0.1), ch_t - Inches(0.1),
        ch_w + Inches(0.2), ch_h + Inches(0.5), RGBColor(0x0F, 0x14, 0x26))

for spd in [0, 20, 40, 60, 80, 100]:
    y = ch_t + ch_h - int(spd / max_s * ch_h)
    add_rect(s3, ch_l, y, ch_w, Pt(0.7), RGBColor(0x22, 0x2A, 0x3F))
    add_text(s3, Inches(0.3), y - Inches(0.12), Inches(0.9), Inches(0.25),
             f"{spd}", 9, C_GRAY, False, PP_ALIGN.RIGHT)

for t in range(0, 180, 10):
    x = ch_l + int(t / max_t * ch_w)
    if t > 0:
        add_rect(s3, x, ch_t, Pt(0.5), ch_h, RGBColor(0x1A, 0x22, 0x35))
    add_text(s3, x - Inches(0.2), ch_t + ch_h + Inches(0.05), Inches(0.4), Inches(0.25),
             f"{t}s", 8, C_GRAY, False, PP_ALIGN.CENTER)

phase_zones = [
    (0, 26, P_INIT, "ECU INIT"), (26, 46, P_WSS, "WSS"),
    (46, 101, P_SPEED, "SPEED & CRUISE"), (101, 126, P_BRAKE, "BRAKE & ABS"),
    (126, 149, P_PARK, "REVERSE & PARKING"), (149, 167, P_FIN, "FINISH"),
]
for t1, t2, clr, label in phase_zones:
    x1 = ch_l + int(t1 / max_t * ch_w)
    x2 = ch_l + int(t2 / max_t * ch_w)
    bar = add_box(s3, x1, ch_t - Inches(0.35), x2 - x1, Inches(0.25), clr)
    box_text(bar, label, 7, C_WHITE, True)

timeline = [
    (0,0),(3,0),(6,0),(11,0),(14,0),(17,0),(20,0),(23,0),(26,0),
    (30,5),(40,5),(43,0),(46,0),
    (56,40),(61,40),(71,80),(91,80),(101,100),
    (111,90),(114,80),(117,60),(119,55),(123,10),
    (126,0),(131,0),(136,5),(139,0),
    (144,0),(149,0),(152,0),(156,0),(159,0),(162,0),(167,0)
]

for i in range(len(timeline) - 1):
    t1, s1 = timeline[i]
    t2, s2 = timeline[i + 1]
    x1 = ch_l + int(t1 / max_t * ch_w)
    y1 = ch_t + ch_h - int(s1 / max_s * ch_h)
    x2 = ch_l + int(t2 / max_t * ch_w)
    y2 = ch_t + ch_h - int(s2 / max_s * ch_h)

    if max(s1, s2) >= 60: clr = P_BRAKE
    elif max(s1, s2) >= 30: clr = P_SPEED
    elif max(s1, s2) > 0: clr = P_WSS
    else: clr = RGBColor(0x44, 0x55, 0x66)

    if s1 == s2:
        add_rect(s3, x1, y1 - Pt(1.5), x2 - x1, Pt(3), clr)
    else:
        add_rect(s3, x1, y1 - Pt(1.5), x2 - x1, Pt(3), clr)
        if i + 1 < len(timeline) - 1:
            vy_top = min(y1, y2)
            vy_bot = max(y1, y2)
            add_rect(s3, x2 - Pt(1.5), vy_top, Pt(3), vy_bot - vy_top, clr)

markers = [
    (35, 5, "WSS 5km/h", P_WSS),
    (58, 40, "SMT 40km/h", P_SPEED),
    (81, 80, "CRUISE 80km/h", P_SPEED),
    (101, 100, "MAX 100km/h", P_BRAKE),
    (114, 80, "DRAG 80km/h", P_BRAKE),
    (117, 60, "BRAKE 60km/h", P_BRAKE),
    (121, 10, "ABS DYNAMIC", P_BRAKE),
    (136, 5, "REVERSE 5km/h", P_PARK),
]
for t, spd, label, clr in markers:
    x = ch_l + int(t / max_t * ch_w)
    y = ch_t + ch_h - int(spd / max_s * ch_h)
    dot_sz = Pt(10)
    dot = s3.shapes.add_shape(MSO_SHAPE.OVAL, x - dot_sz//2, y - dot_sz//2, dot_sz, dot_sz)
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    lbl_y = y - Inches(0.35) if spd > 15 else y - Inches(0.3)
    add_text(s3, x - Inches(0.5), lbl_y, Inches(1.0), Inches(0.3), label, 7, clr, True, PP_ALIGN.CENTER)

add_text(s3, Inches(0.1), ch_t - Inches(0.1), Inches(1.0), Inches(0.25),
         "km/h", 9, C_GRAY, True, PP_ALIGN.RIGHT)

# ================================================================
# Slide 4: Phase Detail (with SYSTEM column)
# ================================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s4, BG_DARK)

add_text(s4, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5), "PHASE DETAIL", 24, C_WHITE, True)
add_text(s4, Inches(0.5), Inches(0.7), Inches(12), Inches(0.3),
         "Step-by-step breakdown with timing, speed, system, and description", 11, C_GRAY)

# System legend (top right)
for i, (label, clr) in enumerate([("ECU", S_ECU), ("PLC", S_PLC), ("NI", S_NI)]):
    lx = Inches(13.5) + i * Inches(0.7)
    lb = add_box(s4, lx, Inches(0.4), Inches(0.55), Inches(0.25), clr)
    box_text(lb, label, 8, RGBColor(0,0,0), True)

card_w = Inches(4.9)
card_h = Inches(3.3)
card_gap_x = Inches(0.2)
card_gap_y = Inches(0.2)

# num, name, dur, spd, systems, note
phase_details = [
    ("PHASE 1", "ECU INITIALIZE", "0-26s  |  26s", P_INIT, [
        ("01","INITIALIZATION","3s","-","P","System boot"),
        ("02","ECU CONNECT","3s","-","E","Diagnostic link"),
        ("03","START COMM","5s","-","E","UDS session start"),
        ("04","ECU IDENT","3s","-","E","ECU version read"),
        ("05","CHECK SIGNAL","3s","-","E","Brake switch check"),
        ("06","DTC READ","3s","-","E","Read fault codes"),
        ("07","DTC CLEAR","3s","-","E","Clear fault codes"),
        ("08","[N] STOP","3s","-","P","Neutral stop"),
    ]),
    ("PHASE 2", "WSS TEST", "26-46s  |  20s", P_WSS, [
        ("09","WSS READY","4s","5","PN","Roll start"),
        ("10","WSS TEST","10s","5","EPN","4-wheel speed check"),
        ("11","[N] STOP","3s","-","P","Neutral stop"),
        ("12","MODE CHECK","3s","-","E","Drive mode verify"),
    ]),
    ("PHASE 3", "SPEED & CRUISE", "46-101s  |  55s", P_SPEED, [
        ("13","SPEED UP","10s","40","PN","Accelerate to 40"),
        ("14","SMT TEST","5s","40","EPN","Speedometer check"),
        ("15","CRUISE UP","10s","80","PN","Accelerate to 80"),
        ("16","CRUISE CHECK","20s","80","EPN","Cruise control test"),
        ("17","MAX SPEED","10s","100","PN","Max speed test"),
    ]),
    ("PHASE 4", "BRAKE & ABS", "101-126s  |  25s", P_BRAKE, [
        ("18","ENGINE BRAKE","10s","90","PN","Engine braking"),
        ("19","DRAG TEST","3s","80","PN","Rolling resistance"),
        ("20","BRAKE TEST","3s","60","PN","Base brake force"),
        ("21","DYNAMIC READY","2s","55","PN","Pedal release"),
        ("22","DYNAMIC TEST","4s","10","EPN","ABS reduction/recovery"),
        ("23","[N] STOP","3s","-","P","Neutral stop"),
    ]),
    ("PHASE 5", "REVERSE & PARKING", "126-149s  |  23s", P_PARK, [
        ("24","REVERSE [R]","5s","-","P","Shift to reverse"),
        ("25","[R] 5KPH","5s","5","PN","Reverse speed check"),
        ("26","[N] STOP","3s","-","P","Neutral stop"),
        ("27","PARKING READY","5s","-","P","Parking brake pull"),
        ("28","PARKING TEST","5s","-","PN","Parking brake hold"),
    ]),
    ("PHASE 6", "FINISH", "149-167s  |  18s", P_FIN, [
        ("29","IDLE","3s","-","E","Idle stability"),
        ("30","DTC READ","4s","-","E","Final fault check"),
        ("31","DTC CLEAR","3s","-","E","Clear fault codes"),
        ("32","ECU DISCONNECT","3s","-","E","Close session"),
        ("33","FINISH","5s","-","P","Test complete"),
    ]),
]

for pi, (pnum, pname, tinfo, color, steps) in enumerate(phase_details):
    row = pi // 3
    col = pi % 3
    cx = Inches(0.4) + col * (card_w + card_gap_x)
    cy = Inches(1.2) + row * (card_h + card_gap_y)

    add_box(s4, cx, cy, card_w, card_h, BG_CARD)
    add_box(s4, cx, cy, card_w, Inches(0.55), color)

    add_text(s4, cx + Inches(0.15), cy + Inches(0.03), Inches(1.0), Inches(0.2),
             pnum, 9, RGBColor(0xDD, 0xDD, 0xDD))
    add_text(s4, cx + Inches(0.15), cy + Inches(0.22), Inches(2.5), Inches(0.25),
             pname, 13, C_WHITE, True)
    add_text(s4, cx + card_w - Inches(2.2), cy + Inches(0.18), Inches(2.0), Inches(0.25),
             tinfo, 9, RGBColor(0xDD, 0xDD, 0xDD), False, PP_ALIGN.RIGHT)

    # Table header
    th_y = cy + Inches(0.65)
    add_text(s4, cx+Inches(0.1), th_y, Inches(0.25), Inches(0.2), "#", 7, C_GRAY, True)
    add_text(s4, cx+Inches(0.35), th_y, Inches(1.3), Inches(0.2), "STEP", 7, C_GRAY, True)
    add_text(s4, cx+Inches(1.7), th_y, Inches(0.45), Inches(0.2), "TIME", 7, C_GRAY, True)
    add_text(s4, cx+Inches(2.15), th_y, Inches(0.5), Inches(0.2), "km/h", 7, C_GRAY, True)
    add_text(s4, cx+Inches(2.65), th_y, Inches(0.9), Inches(0.2), "SYSTEM", 7, C_GRAY, True)
    add_text(s4, cx+Inches(3.6), th_y, Inches(1.2), Inches(0.2), "NOTE", 7, C_GRAY, True)

    add_rect(s4, cx + Inches(0.1), th_y + Inches(0.2),
             card_w - Inches(0.2), Pt(0.7), RGBColor(0x33, 0x3D, 0x55))

    row_h = Inches(0.27)
    for si, (num, name, dur, spd, sys, note) in enumerate(steps):
        ry = th_y + Inches(0.25) + si * row_h
        tc = C_WHITE if si % 2 == 0 else RGBColor(0xCC, 0xCC, 0xCC)
        add_text(s4, cx+Inches(0.1), ry, Inches(0.25), row_h, num, 8, color, True)
        add_text(s4, cx+Inches(0.35), ry, Inches(1.3), row_h, name, 8, tc)
        add_text(s4, cx+Inches(1.7), ry, Inches(0.45), row_h, dur, 8, tc)
        sc = P_SPEED if spd != "-" else C_GRAY
        add_text(s4, cx+Inches(2.15), ry, Inches(0.5), row_h, spd, 8, sc)
        # System badges
        badges = get_sys_badges(sys)
        add_sys_badges(s4, cx + Inches(2.65), ry + Inches(0.06), badges)
        add_text(s4, cx+Inches(3.6), ry, Inches(1.2), row_h, note, 7, C_GRAY)

output = r"E:\파키스탄 LX3\KI-RnB_1208\LX3_ABS_Test_Sequence.pptx"
prs.save(output)
print(f"Done: {output}")
